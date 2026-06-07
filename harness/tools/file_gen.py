from __future__ import annotations

import re
import time
import uuid
from pathlib import Path
from typing import Any

from harness import config


def _generated_dir() -> Path:
    directory = Path(getattr(config, "GENERATED_DIR"))
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def _safe_stem(title: str, fallback: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9가-힣_-]+", "_", str(title or "").strip()).strip("_")
    return (stem or fallback)[:60]


def _unique_path(title: str, suffix: str, fallback: str) -> Path:
    stamp = time.strftime("%Y%m%d_%H%M%S")
    token = uuid.uuid4().hex[:8]
    return _generated_dir() / f"{_safe_stem(title, fallback)}_{stamp}_{token}{suffix}"


def _string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value]
    if value is None:
        return []
    return [str(value)]


def create_pptx(title: str, slides: list) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "error: 이 기능은 python-pptx 설치가 필요합니다."

    path = _unique_path(title, ".pptx", "presentation")
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = str(title)
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = ""

    for item in slides or []:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title", ""))
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = _string_list(item.get("bullets", []))
        if bullets:
            body.text = bullets[0]
            for bullet in bullets[1:]:
                paragraph = body.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
        notes = item.get("notes")
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = str(notes)
            except Exception:
                pass

    presentation.save(str(path))
    return f"생성 완료: {path}"


def create_docx(title: str, sections: list) -> str:
    try:
        from docx import Document
    except ImportError:
        return "error: 이 기능은 python-docx 설치가 필요합니다."

    path = _unique_path(title, ".docx", "document")
    document = Document()
    document.add_heading(str(title), 0)

    for section in sections or []:
        heading = section.get("heading")
        body = section.get("body", "")
        if heading:
            document.add_heading(str(heading), level=1)
        for paragraph in str(body).splitlines():
            if paragraph.strip():
                document.add_paragraph(paragraph.strip())

    document.save(str(path))
    return f"생성 완료: {path}"


def create_xlsx(title: str, sheets: list) -> str:
    try:
        import openpyxl
        from openpyxl.styles import Font
    except ImportError:
        return "error: 이 기능은 openpyxl 설치가 필요합니다."

    path = _unique_path(title, ".xlsx", "workbook")
    workbook = openpyxl.Workbook()
    default = workbook.active
    workbook.remove(default)

    for index, sheet_data in enumerate(sheets or [], start=1):
        raw_name = str(sheet_data.get("name") or f"Sheet{index}")
        name = re.sub(r"[\[\]:*?/\\]", "_", raw_name)[:31] or f"Sheet{index}"
        sheet = workbook.create_sheet(title=name)
        headers = _string_list(sheet_data.get("headers", []))
        if headers:
            sheet.append(headers)
            for cell in sheet[1]:
                cell.font = Font(bold=True)
        for row in sheet_data.get("rows", []) or []:
            if isinstance(row, list):
                sheet.append(row)
            else:
                sheet.append([row])

    if not workbook.worksheets:
        workbook.create_sheet(title="Sheet1")

    workbook.save(str(path))
    return f"생성 완료: {path}"


def create_pdf(title: str, html: str) -> str:
    try:
        from weasyprint import HTML
    except ImportError:
        return "error: 이 기능은 weasyprint 설치가 필요합니다."

    path = _unique_path(title, ".pdf", "document")
    HTML(string=str(html)).write_pdf(str(path))
    return f"생성 완료: {path}"


def register(registry: Any) -> None:
    registry.add(
        "create_pptx",
        "슬라이드 목록으로 PPTX 파일을 생성합니다.",
        {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "notes": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["title", "slides"],
        },
        create_pptx,
        permissions=(),
    )
    registry.add(
        "create_docx",
        "섹션 목록으로 DOCX 파일을 생성합니다.",
        {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "sections": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "body": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["title", "sections"],
        },
        create_docx,
        permissions=(),
    )
    registry.add(
        "create_xlsx",
        "시트 데이터로 XLSX 파일을 생성합니다.",
        {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "sheets": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string"},
                            "headers": {"type": "array", "items": {"type": "string"}},
                            "rows": {"type": "array", "items": {"type": "array"}},
                        },
                    },
                },
            },
            "required": ["title", "sheets"],
        },
        create_xlsx,
        permissions=(),
    )
    registry.add(
        "create_pdf",
        "HTML을 PDF 파일로 렌더링합니다.",
        {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "html": {"type": "string"},
            },
            "required": ["title", "html"],
        },
        create_pdf,
        permissions=(),
    )

